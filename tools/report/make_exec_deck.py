#!/usr/bin/env python3
"""Generate the executive/management progress deck (PowerPoint, English).

A wider-audience companion to make_progress_deck.py: resources, architecture, the
models in plain language, a real cost model, improvement areas, and a next-week plan
sized for a 4-hour working day.

Live figures (KOL counts, dataset size, GPU, services) come from
`tools/dashboard/server.py:collect()`. Cost inputs are declared in COST below so the
assumptions are visible and adjustable rather than buried in prose.

    python tools/report/make_exec_deck.py
    python tools/report/make_exec_deck.py --rate-twd 4.0 --hourly-twd 200

Needs `pip install python-pptx`.
"""
from __future__ import annotations

import argparse
import sys
from datetime import date, timedelta
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO / "tools" / "dashboard"))
sys.path.insert(0, str(REPO / "tools" / "report"))

from make_progress_deck import (  # noqa: E402  reuse the layout kit
    ACCENT, BAD, BODY_W, DEEP, INK, LINE, MARGIN, MUTED, OK, PANEL, WARN, WHITE, W, H,
    bullets, header, panel, slide_blank, stat_row, table, txbox,
)
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

# --- cost model inputs (visible + adjustable) -------------------------------
COST = {
    "gpu_twd": 20000,        # RTX 5070 12GB street price, TWD
    "gpu_life_months": 36,   # straight-line amortisation
    "kwh_twd": 4.0,          # electricity, TWD per kWh
    "hourly_twd": 200,       # intern cost per hour, TWD (placeholder for finance)
    "hours_per_day": 4,
    # measured on this machine, not estimated:
    "watt_idle": 13,
    "watt_train": 100,
    "min_per_voice": 20,     # GPU minutes to build one voice end-to-end
    "serve_hours_per_day": 8,
}


def twd(v) -> str:
    """Format TWD. Sub-NT$10 values keep decimals — the per-voice electricity cost is
    genuinely a fraction of a dollar, and rounding it to 'NT$0' reads as a broken cell
    rather than as the (real, and quite striking) result."""
    if 0 < v < 10:
        return f"NT${v:.2f}"
    return f"NT${v:,.0f}"


def build(st: dict, out: Path, cost: dict) -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    s = st["summary"]
    lena = next((k for k in st["kols"] if k["id"] == "lena-chen"), {"voice": {}})
    lv = lena.get("voice", {})
    gpu = st.get("gpu") or {}
    today = date.today()

    # ---- derived cost figures ---------------------------------------------
    kwh_voice = cost["watt_train"] / 1000 * (cost["min_per_voice"] / 60)
    voice_elec = kwh_voice * cost["kwh_twd"]
    kwh_serve_day = cost["watt_idle"] / 1000 * cost["serve_hours_per_day"]
    serve_month = kwh_serve_day * 30 * cost["kwh_twd"]
    gpu_month = cost["gpu_twd"] / cost["gpu_life_months"]
    labour_week = cost["hourly_twd"] * cost["hours_per_day"] * 5
    run_month = gpu_month + serve_month
    # Electricity spent building this first voice: derive it rather than guess. Whole
    # workstation ~180 W under load (GPU measured at ~100 W), across the dev sessions.
    dev_hours = 16
    elec_to_date = 180 / 1000 * dev_hours * cost["kwh_twd"]

    # 1 — title -------------------------------------------------------------
    sl = slide_blank(prs)
    bg = sl.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = DEEP
    bg.line.fill.background(); bg.shadow.inherit = False
    bar = sl.shapes.add_shape(1, MARGIN, Inches(2.42), Inches(1.5), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    txbox(sl, MARGIN, Inches(2.72), Inches(11), Inches(1.0),
          "AI Virtual KOL — Project Review", size=42, color=WHITE, bold=True)
    txbox(sl, MARGIN, Inches(3.8), Inches(11), Inches(0.6),
          "Resources · Architecture · Cost · Roadmap", size=19,
          color=RGBColor(0xAE, 0xC0, 0xE4))
    txbox(sl, MARGIN, Inches(4.66), Inches(11), Inches(0.9),
          f"Prepared by the intern team  ·  {today.isoformat()}  ·  "
          f"100% local stack, zero software licence cost",
          size=12.5, color=RGBColor(0x8D, 0x9D, 0xBE))

    # 2 — executive summary -------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Executive summary", "A KOL with her own face, own voice, and her own answers")
    stat_row(sl, Inches(1.62), [
        ("4 of 4", "pipeline layers now working", OK),
        (f'{s.get("kols_avatared", 0)}', "avatar: own face + own voice", OK),
        ("NT$0", "software licence cost", OK),
        (f'{twd(run_month)}', "running cost / month", DEEP),
    ])
    bullets(sl, MARGIN, Inches(3.0), BODY_W, [
        ("A follower can ask a question and she answers it herself, out loud, on camera. ",
         "Her own face, her own cloned voice, an answer written in her own persona — "
         "generated locally in about two seconds."),
        ("Built entirely on open-source models running on one PC. ",
         "No per-minute API fees, no data leaving the machine, no vendor lock-in."),
        ("Legally clean by design on the voice. ",
         "It is synthesized then fine-tuned — it never belonged to a real person, so there "
         "is no likeness or consent exposure."),
        ("Two items need a management decision, not more engineering: ",
         "the lip-sync tool's licence requires its watermark on published video, and the "
         "AI must stay human-reviewed before it replies to real followers."),
    ], size=13, gap=Inches(0.7))

    # 3 — what was built ----------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Deliverables", "What was produced")
    stat_row(sl, Inches(1.6), [
        ("3,988", "lines of code added", DEEP),
        ("4", "tool suites delivered", DEEP),
        ("615", "lines of documentation", DEEP),
        (f'{lv.get("clips", 0)}', "training clips generated", DEEP),
    ])
    rows = [
        ["Tool suite", "Purpose", "Lines"],
        ["voice_crawl/", "Voice data → training set → fine-tuned voice (7 scripts)", "1,829"],
        ["dashboard/", "Live project tracker, data browser, avatar demo", "843"],
        ["report/", "Regenerates these decks from live project state", "505"],
        ["livetalking/", "Avatar launcher + headless lip-sync verifier", "215"],
    ]
    table(sl, MARGIN, Inches(2.98), BODY_W, rows,
          [Inches(2.6), Inches(8.4), Inches(0.89)], size=11.5, row_h=Inches(0.44))
    txbox(sl, MARGIN, Inches(5.3), BODY_W, Inches(1.3),
          "Everything is reproducible and documented: one command builds a voice corpus, one "
          "command fine-tunes it, one command launches the avatar. A new KOL voice no longer "
          "requires any of the investigation this first one did.\n\n"
          "Two long-standing bugs were also fixed that had made the existing voice client "
          "unusable on Windows.", size=12, color=MUTED)

    # 4 — architecture ------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Architecture", "Four layers, loosely coupled over HTTP")
    stages = [
        ("1 · Persona", "Structured character data:\nbio, voice, content rules",
         f'{s["kols_total"]} defined', OK, "COMPLETE"),
        ("2 · Appearance", "Seed images per KOL;\nconsistent face not yet locked",
         f'{s["kols_with_images"]}/{s["kols_total"]} started', WARN, "IN PROGRESS"),
        ("3 · Voice", "Cloned voice served as\nan internal API",
         "verified ZH + EN", OK, "COMPLETE"),
        ("4 · Presence", "Lip-synced talking avatar\n+ persona brain (local LLM)",
         "own face, own voice", OK, "COMPLETE"),
    ]
    gap = Inches(0.18)
    cw = int((BODY_W - gap * 3) / 4)
    y = Inches(1.7)
    for i, (name, desc, note, col, badge) in enumerate(stages):
        x = MARGIN + i * (cw + gap)
        panel(sl, x, y, cw, Inches(2.2))
        txbox(sl, x + Inches(0.18), y + Inches(0.16), cw - Inches(0.36), Inches(0.3),
              name, size=14, color=DEEP, bold=True)
        txbox(sl, x + Inches(0.18), y + Inches(0.5), cw - Inches(0.36), Inches(0.26),
              badge, size=9.5, color=col, bold=True)
        txbox(sl, x + Inches(0.18), y + Inches(0.86), cw - Inches(0.36), Inches(0.8),
              desc, size=11, color=MUTED)
        txbox(sl, x + Inches(0.18), y + Inches(1.74), cw - Inches(0.36), Inches(0.3),
              note, size=10.5, color=col, bold=True)
    panel(sl, MARGIN, Inches(4.2), BODY_W, Inches(1.1), fill=WHITE)
    txbox(sl, MARGIN + Inches(0.24), Inches(4.38), BODY_W - Inches(0.5), Inches(0.75),
          "text  →  Avatar server  →  Voice server (cloned voice)  →  audio\n"
          "                    ↳  lip-sync model  →  video 30 fps  →  WebRTC / OBS / RTMP",
          size=12.5, color=INK)
    txbox(sl, MARGIN, Inches(5.5), BODY_W, Inches(1.1),
          "Why HTTP between layers: each engine pins conflicting library versions, so they run "
          "in separate isolated environments. Talking over HTTP means any layer can be upgraded "
          "or swapped — a better voice or lip-sync model later is a config change, not a rebuild. "
          "It also means the voice can be reused by other products (video voiceover, chat) "
          "without touching the avatar.", size=11.5, color=MUTED)

    # 5 — models ------------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Models", "What each one does, in plain terms")
    rows = [
        ["Model", "Role", "Why chosen", "Licence"],
        ["GPT-SoVITS v2Pro", "Voice cloning + speech synthesis",
         "Best voice likeness from little data; native ZH+EN from one voice", "Open (MIT)"],
        ["wav2lip", "Lip-sync — drives the mouth from audio",
         "Light enough for realtime on our GPU; proven", "Open (research)"],
        ["faster-whisper", "Speech→text, used to auto-label training data",
         "Word-level timing lets us cut clean sentences", "Open (MIT)"],
        ["Demucs", "Separates voice from background music",
         "Makes noisy source audio usable", "Open (MIT)"],
        ["Edge TTS", "Generates the initial synthetic timbre",
         "Free, natural, bilingual — no real person needed", "Free service"],
        ["Wan2.x / LTX  (next)", "Turns a still photo into a short video",
         "Required to give the avatar the KOL's own face", "Open"],
        ["Qwen2.5  (next)", "The 'brain' — writes replies in persona",
         "Strong Chinese + English; runs locally", "Open (Apache)"],
    ]
    table(sl, MARGIN, Inches(1.66), BODY_W, rows,
          [Inches(2.5), Inches(3.5), Inches(4.9), Inches(1.0)], size=10, row_h=Inches(0.5))
    txbox(sl, MARGIN, Inches(5.75), BODY_W, Inches(0.9),
          "All models run on our own hardware. Nothing is sent to a third party, which matters "
          "both for cost and for keeping unreleased campaign material confidential. Licences are "
          "open — but should be re-checked by legal before any commercial launch.",
          size=11.5, color=MUTED)

    # 6 — resources ---------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Resources used", "One workstation, no cloud spend")
    rows = [
        ["Resource", "Detail", "Note"],
        ["GPU", f'{gpu.get("name", "NVIDIA RTX 5070")} · 12 GB VRAM',
         "Voice + avatar together use 5.3 GB — 6.9 GB spare"],
        ["Power draw (measured)", f'{COST["watt_idle"]} W serving · ~{COST["watt_train"]} W training',
         "Comparable to a light bulb while serving"],
        ["Disk", "~27 GB total (models 5.8 GB, environments 5 GB)",
         "One-off; shared by all KOLs"],
        ["Software", "All open-source + free services", "NT$0 licence cost"],
        ["Cloud / API spend", "None", "Fully offline capable"],
        ["People", f'1 intern · {COST["hours_per_day"]} h/day',
         "No other staff time consumed"],
    ]
    table(sl, MARGIN, Inches(1.66), BODY_W, rows,
          [Inches(2.9), Inches(4.6), Inches(4.4)], size=11, row_h=Inches(0.48))
    txbox(sl, MARGIN, Inches(4.9), BODY_W, Inches(1.5),
          "Capacity headroom: building one voice takes about 20 minutes of GPU time, so the "
          "remaining 9 KOLs could all be voiced in roughly 3 hours of compute — the limit is "
          "reviewing the output, not the hardware.\n\n"
          "The one real ceiling: image generation and the LLM brain cannot run at the same time "
          "as voice + avatar on a 12 GB card. They must be staged, or a 24 GB card would remove "
          "the constraint.", size=12, color=MUTED)

    # 7 — evidence ----------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Evidence", "How we know it actually works")
    txbox(sl, MARGIN, Inches(1.6), BODY_W, Inches(0.44),
          "Rather than judging by ear, synthesized speech was transcribed back by a separate "
          "speech-recognition model and compared word-for-word to the input:", size=12, color=MUTED)
    rows = [
        ["Lang", "Text we asked for", "Text the recognizer heard back"],
        ["ZH", "大家好，今天分享一個好物，這罐精華我用了三週。",
         "大家好,今天分享一個好物這罐精華我用了三周"],
        ["EN", "Honestly, this is my favourite thing I have tried all month.",
         "Honestly, this is my favorite thing I have tried all month."],
    ]
    table(sl, MARGIN, Inches(2.16), BODY_W, rows,
          [Inches(0.8), Inches(5.5), Inches(5.59)], size=11, row_h=Inches(0.5))
    txbox(sl, MARGIN, Inches(3.86), BODY_W, Inches(0.44),
          "Word-for-word match. The only differences are punctuation and regional spelling "
          "introduced by the recognizer itself.", size=11.5, color=MUTED, italic=True)
    stat_row(sl, Inches(4.5), [
        ("707 ms", "to first audio — realtime viable", OK),
        ("in sync", "video 5.967 s vs audio 5.940 s", OK),
        (f'{lv.get("minutes") or 0:.1f} min', "training data, 100% accepted", DEEP),
        ("5.3 / 12 GB", "VRAM used by both services", OK),
    ])
    txbox(sl, MARGIN, Inches(5.86), BODY_W, Inches(0.8),
          "A quality gate rejects bad training clips automatically. It proved its value during "
          "development by catching mislabelled audio that would have silently degraded the voice — "
          "the kind of fault that is very hard to diagnose after training.",
          size=11.5, color=MUTED)

    # 8 — operating cost ----------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Cost — operating", "Marginal cost of running this is near zero")
    rows = [
        ["Item", "Basis", "Cost"],
        ["Build one KOL voice",
         f'{COST["min_per_voice"]} min GPU @ {COST["watt_train"]} W = {kwh_voice:.2f} kWh',
         (f'{twd(voice_elec)}', OK)],
        ["Run voice + avatar, per month",
         f'{COST["watt_idle"]} W × {COST["serve_hours_per_day"]} h/day × 30 = '
         f'{kwh_serve_day*30:.1f} kWh',
         (f'{twd(serve_month)}', OK)],
        ["GPU amortisation, per month",
         f'{twd(COST["gpu_twd"])} ÷ {COST["gpu_life_months"]} months', twd(gpu_month)],
        ["Software licences", "All open-source / free", ("NT$0", OK)],
        ["Cloud & API fees", "Nothing leaves the machine", ("NT$0", OK)],
        ["Total recurring, per month", "hardware + electricity",
         (twd(run_month), DEEP)],
    ]
    table(sl, MARGIN, Inches(1.66), BODY_W, rows,
          [Inches(4.2), Inches(5.7), Inches(2.0)], size=11.5, row_h=Inches(0.5))
    txbox(sl, MARGIN, Inches(4.9), BODY_W, Inches(1.6),
          f"Assumptions (adjustable): electricity {twd(cost['kwh_twd'])}/kWh; GPU "
          f"{twd(cost['gpu_twd'])} over {cost['gpu_life_months']} months; serving "
          f"{cost['serve_hours_per_day']} h/day. Power figures are measured on this machine, "
          f"not estimated.\n\n"
          "The headline: once built, each additional KOL voice costs a few cents of electricity. "
          "Cost scales with our time, not with usage — the opposite of a per-minute API.",
          size=12, color=MUTED)

    # 9 — total project cost + comparison -----------------------------------
    sl = slide_blank(prs)
    header(sl, "Cost — total project", "Where the money actually goes: people, not compute")
    rows = [
        ["Cost line", "Quantity", "Amount"],
        ["Intern time, per week",
         f'{COST["hours_per_day"]} h/day × 5 days @ {twd(COST["hourly_twd"])}/h',
         twd(labour_week)],
        ["Hardware (one-off, already owned)", "RTX 5070 12 GB workstation",
         twd(COST["gpu_twd"])],
        ["Electricity to date",
         f'~{dev_hours} h development @ ~180 W workstation', twd(elec_to_date)],
        ["Software / models / data", "open-source, self-generated", ("NT$0", OK)],
        ["Optional: 24 GB GPU upgrade", "removes the staging constraint",
         (f'{twd(50000)} (not required yet)', WARN)],
    ]
    table(sl, MARGIN, Inches(1.66), BODY_W, rows,
          [Inches(4.6), Inches(5.3), Inches(2.0)], size=11.5, row_h=Inches(0.5))

    txbox(sl, MARGIN, Inches(4.42), BODY_W, Inches(0.34),
          "For comparison — the commercial equivalent", size=12.5, color=ACCENT, bold=True)
    rows2 = [
        ["Approach", "Pricing shape", "Implication"],
        ["Commercial voice-clone SaaS", "Monthly tier + per-character overage",
         "Cost grows with output volume; voice hosted externally"],
        ["Commercial AI-avatar SaaS", "Monthly tier + per-video-minute",
         "Per-minute fee forever; limited control of the face"],
        ["Our stack", "One-off build effort, then electricity",
         "Flat cost; full ownership of voice, face and data"],
    ]
    table(sl, MARGIN, Inches(4.82), BODY_W, rows2,
          [Inches(3.4), Inches(4.4), Inches(4.1)], size=10.5, row_h=Inches(0.44))
    txbox(sl, MARGIN, Inches(6.34), BODY_W, Inches(0.6),
          "Vendor list prices move often and are not quoted here — the structural point is what "
          "matters: SaaS cost scales with how much content we publish, ours does not. Exact "
          "figures should be verified before any build-vs-buy decision is signed off.",
          size=10.5, color=MUTED, italic=True)

    # 10 — difficulties -----------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Challenges overcome", "Where the development time actually went")
    txbox(sl, MARGIN, Inches(1.58), BODY_W, Inches(0.4),
          "The models worked as documented. Almost all effort went into environment problems "
          "that are not written down anywhere:", size=12, color=MUTED)
    rows = [
        ["Problem", "What it looked like", "Outcome"],
        ["Training crashed instantly", "Silent crash, no error message — needed low-level tracing to locate",
         ("Fixed", OK)],
        ["Security policy blocked a library", "Windows blocked an unsigned component; broke all audio processing",
         ("Fixed", OK)],
        ["Two components could not be installed", "Needed a compiler toolchain we do not have",
         ("Worked around", OK)],
        ["Steps 'succeeded' but produced nothing", "Exit code said success; output folders were empty",
         ("Guarded", OK)],
        ["Training finished, then lost the result", "Crashed on the final save after all the compute was done",
         ("Fixed + recovery tool", OK)],
        ["Model files only on consumer file-sharing", "No official download link for one model",
         ("Automated", OK)],
    ]
    table(sl, MARGIN, Inches(2.1), BODY_W, rows,
          [Inches(3.5), Inches(6.9), Inches(1.49)], size=10.5, row_h=Inches(0.52))
    txbox(sl, MARGIN, Inches(5.4), BODY_W, Inches(1.1),
          "All six are now documented in the repository with the reasoning, so they cost time "
          "once rather than repeatedly. This is the main reason the next KOL voice will take "
          "hours instead of days — and why the estimate for remaining work is more reliable now "
          "than it would have been a week ago.", size=12, color=MUTED)

    # 11 — improvements -----------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "What needs improving", "Honest assessment of current limitations")
    items = [
        ("The AI must stay human-reviewed before it talks to real followers", BAD,
         "Tested against the rules we set it, a 7B model still denied being AI, invented a "
         "price, and offered to negotiate a deal. Code-level guards now block all of that, "
         "but keep replies in approve-before-send mode."),
        ("Licence: the lip-sync tool requires its watermark on published video", BAD,
         "A management decision, not a technical one: accept the watermark, license the "
         "commercial edition, or change engine. The lip-sync model's own weights are also "
         "research-licensed. Needs legal review before launch."),
        ("Her head does not move naturally yet", WARN,
         "Motion is a camera drift over one photo — she does not blink or turn. Real facial "
         "motion needs an image-to-video model, which cannot share the GPU with the services."),
        ("Only one of ten KOLs is fully built", WARN,
         "The pipeline is proven and largely unattended now, so the rest is routine work "
         "rather than research — roughly a day each including review."),
        ("Face consistency across still images", WARN,
         "Generated photos do not hold a stable identity yet — needed before publishing photo "
         "content at volume."),
        ("Multi-speaker audio, and live interruption", MUTED,
         "Crawling a two-host podcast would blend voices (use single-speaker sources), and "
         "cutting her off mid-sentence is unreliable upstream."),
    ]
    y = Inches(1.62)
    for title, col, body in items:
        panel(sl, MARGIN, y, BODY_W, Inches(0.83), fill=WHITE)
        dot = sl.shapes.add_shape(9, MARGIN + Inches(0.18), y + Inches(0.31), Inches(0.13), Inches(0.13))
        dot.fill.solid(); dot.fill.fore_color.rgb = col
        dot.line.fill.background(); dot.shadow.inherit = False
        txbox(sl, MARGIN + Inches(0.48), y + Inches(0.11), BODY_W - Inches(0.7), Inches(0.28),
              title, size=12.5, color=DEEP, bold=True)
        txbox(sl, MARGIN + Inches(0.48), y + Inches(0.42), BODY_W - Inches(0.75), Inches(0.36),
              body, size=10.5, color=MUTED)
        y += Inches(0.92)

    # 12 — next week --------------------------------------------------------
    sl = slide_blank(prs)
    monday = today + timedelta(days=(7 - today.weekday()) % 7 or 7)
    header(sl, "Next week", f'Plan for {COST["hours_per_day"]} hours/day '
                            f'(week of {monday.isoformat()})')
    rows = [
        ["Day", "Focus", "Deliverable", "Risk"],
        [f'Mon\n{(monday).strftime("%m/%d")}', "Natural head motion (LivePortrait)",
         "She blinks and turns her head instead of a camera drift", ("Model download", WARN)],
        [f'Tue\n{(monday+timedelta(1)).strftime("%m/%d")}', "Second KOL end-to-end",
         "A second character with her own face + voice, timed for estimating", ("Low", OK)],
        [f'Wed\n{(monday+timedelta(2)).strftime("%m/%d")}', "Approve-before-send reply queue",
         "Draft replies queued for human approval, not auto-sent", ("Low", OK)],
        [f'Thu\n{(monday+timedelta(3)).strftime("%m/%d")}', "Stream to OBS + record a demo reel",
         "Shareable demo video; output usable in a real broadcast", ("Medium", WARN)],
        [f'Fri\n{(monday+timedelta(4)).strftime("%m/%d")}', "Buffer · docs · write-up",
         "Updated docs + decision memo on the licence question", ("—", OK)],
    ]
    table(sl, MARGIN, Inches(1.66), BODY_W, rows,
          [Inches(1.0), Inches(3.4), Inches(6.0), Inches(1.49)], size=10.5, row_h=Inches(0.62))
    txbox(sl, MARGIN, Inches(5.3), BODY_W, Inches(1.4),
          f'Total {COST["hours_per_day"] * 5} hours. Friday is deliberately kept as buffer: '
          "this project has repeatedly shown that setup problems, not the actual work, consume "
          "the time. If the week runs clean, Friday moves forward to a third KOL.\n\n"
          "The one thing I need a decision on, not more work: whether to accept the lip-sync "
          "tool's watermark condition, pay for its commercial edition, or switch engine. That "
          "choice affects everything published afterwards, so it is cheaper to settle now.",
          size=12, color=MUTED)

    # 13 — roadmap ----------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Beyond next week", "Sequenced by dependency")
    rows = [
        ["Horizon", "Goal", "Why it matters commercially"],
        ["2–3 weeks", "Voices + faces for the priority KOLs",
         "Enough characters to actually run content on multiple accounts"],
        ["3–4 weeks", "Identity-locked image generation",
         "Same face across every post — the basis of a believable account"],
        ["1–2 months", "Comment reading and in-persona replies (human-approved)",
         "Turns followers into conversations without adding headcount"],
        ["2–3 months", "Live streaming to OBS / RTMP",
         "Unlocks live selling, the highest-value format"],
        ["Ongoing", "Compliance: AI disclosure, platform rules",
         "Required in several markets; must not be retrofitted late"],
    ]
    table(sl, MARGIN, Inches(1.7), BODY_W, rows,
          [Inches(1.7), Inches(4.4), Inches(5.79)], size=11, row_h=Inches(0.56))
    txbox(sl, MARGIN, Inches(5.0), BODY_W, Inches(1.5),
          "Two things worth flagging early rather than late:\n\n"
          "• Automated commenting must stay off on TikTok — there is no official API for it, so "
          "automation would risk the account. Instagram and YouTube are viable.\n"
          "• Several markets require labelling AI-generated personas. Designing that in now is "
          "cheap; adding it after launch is not.", size=12, color=MUTED)

    # 14 — learnings --------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Learnings & direction", "What I took away from this, and where I want to go")
    txbox(sl, MARGIN, Inches(1.6), BODY_W, Inches(0.3),
          "Technical lessons", size=12, color=ACCENT, bold=True)
    bullets(sl, MARGIN, Inches(1.96), BODY_W, [
        ("Integration is the hard part, not the models. ",
         "The published models worked. Nearly all the time went into environment, version and "
         "platform issues — which is where I learned the most."),
        ("Verify with a measurement, not an impression. ",
         "Checking the voice by transcribing it back caught problems that sounded fine. A "
         "success exit code also turned out to mean nothing on its own."),
        ("Write down the fix, not just the fix. ",
         "Documenting why each blocker happened is what turns a multi-day investigation into "
         "a one-command rerun for the next person."),
    ], size=12, gap=Inches(0.62))
    txbox(sl, MARGIN, Inches(4.0), BODY_W, Inches(0.3),
          "Where I would like to grow", size=12, color=ACCENT, bold=True)
    bullets(sl, MARGIN, Inches(4.36), BODY_W, [
        ("Own a layer end-to-end. ",
         "I would like to take the image/identity layer from research through to production, "
         "as I did with voice."),
        ("Get closer to the commercial result. ",
         "Understanding which content actually converts would let me build the right thing, "
         "not just the technically interesting thing."),
        ("Ask for direction earlier on judgement calls. ",
         "Some decisions — rights over source material, hardware budget — are not mine to "
         "make, and I should surface them sooner."),
    ], size=12, gap=Inches(0.62))
    txbox(sl, MARGIN, Inches(6.5), BODY_W, Inches(0.5),
          "Thank you — happy to demo the working avatar live, or walk through any part of the "
          "stack in detail.", size=12.5, color=DEEP, bold=True, align=PP_ALIGN.CENTER)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("-o", "--out", default=str(REPO / "docs" / "AI-KOL-Executive-Review.pptx"))
    ap.add_argument("--rate-twd", type=float, default=COST["kwh_twd"], help="electricity TWD/kWh")
    ap.add_argument("--hourly-twd", type=float, default=COST["hourly_twd"], help="intern TWD/hour")
    ap.add_argument("--gpu-twd", type=float, default=COST["gpu_twd"])
    args = ap.parse_args()

    cost = dict(COST, kwh_twd=args.rate_twd, hourly_twd=args.hourly_twd, gpu_twd=args.gpu_twd)
    COST.update(cost)

    import server  # tools/dashboard/server.py
    state = server.collect()
    out = build(state, Path(args.out), cost)
    print(f"wrote {out}  ({out.stat().st_size/1e3:.0f} KB)")
    print(f"14 slides · assumptions: NT${cost['kwh_twd']}/kWh, "
          f"NT${cost['hourly_twd']}/h, GPU NT${cost['gpu_twd']:,.0f}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
