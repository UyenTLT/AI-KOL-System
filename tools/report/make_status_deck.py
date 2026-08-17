#!/usr/bin/env python3
"""Generate the status-review deck (PowerPoint, English) — what works, what does not, what I need.

The third deck, and it answers a different question from the other two. `make_progress_deck.py`
says what was built and `make_exec_deck.py` says what it costs; this one is written for the
review conversation that follows: the two quality problems that are currently blocking a public
launch (replies that read as stiff, and a 7.6 s wait before she speaks), what each one is
actually caused by, and the three things only management can unblock.

Every number is measured on this machine. The sources are named in the comments beside them so a
figure can be traced rather than trusted — the roster and GPU come live from
`tools/dashboard/server.py:collect()`, the rest from the measurement each one was taken in.

    python tools/report/make_status_deck.py
    python tools/report/make_status_deck.py -o docs/AI-KOL-Status-Review.pptx

Needs `pip install python-pptx`.
"""
from __future__ import annotations

import argparse
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO / "tools" / "dashboard"))
sys.path.insert(0, str(REPO / "tools" / "report"))

from make_progress_deck import (  # noqa: E402  reuse the layout kit
    ACCENT, BAD, BODY_W, DEEP, FONT, INK, LINE, MARGIN, MUTED, OK, PANEL, WARN, WHITE, W, H,
    bullets, header, panel, slide_blank, stat_row, table, txbox,
)
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

# --- measured constants -----------------------------------------------------
#
# Declared here rather than inline in the slides so the deck states one number in one place, and
# so the source of each is visible next to it. All are measurements, not estimates.
M = {
    "llm_s": 0.8,          # tools/livestream/server.py — persona reply, tuned model via Ollama
    "tts_s": 6.8,          # same measurement — CosyVoice 2, one conversational reply
    "tuned_pct": 92,       # tools/llm_train/evaluate.py — held-out mid-conversation turns
    "base_pct": 83,        # same run, base model carrying the full 3,645-char persona prompt
    "tuned_prompt": 129,   # chars of system prompt the tuned model needs
    "base_prompt": 3645,
    "train_rows": 303,     # datasets/sofia-vargas-chat-train.jsonl
    "val_rows": 32,
    "tok_gguf": 30.7,      # RUN-TUNED.ps1 — merged q4_K_M through Ollama
    "tok_hf": 13,          # same weights served through transformers 4-bit
    "pitch_range": 14.48,  # tools/studio/README.md — Sofia on CosyVoice 2, 9 sentence types
    "pitch_human": 14.10,  # the real human reference clip
    "pitch_old": 10.43,    # her previous fine-tuned GPT-SoVITS voice — the flat one
    "asr": 98,             # ASR round-trip accuracy, same run
    "spk_cross": 0.691,    # tools/studio/README.md — ERes2NetV2 cross-character mean
    "spk_base": 0.919,     # same-speaker baseline
    "blinks": 14,          # tools/livetalking/README.md — per minute, LivePortrait base clip
    "vram_used": 5287,     # MiB, LiveTalking + api_v2 resident together on a 12 GB card
    "vram_free": 6657,
}
M["turn_s"] = round(M["llm_s"] + M["tts_s"], 1)
M["tts_pct"] = round(M["tts_s"] / M["turn_s"] * 100)


def build(st: dict, out: Path) -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    s = st["summary"]
    gpu = st.get("gpu") or {}
    today = date.today()

    # 1 — title -------------------------------------------------------------
    sl = slide_blank(prs)
    bg = sl.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = DEEP
    bg.line.fill.background(); bg.shadow.inherit = False
    bar = sl.shapes.add_shape(1, MARGIN, Inches(2.42), Inches(1.5), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    txbox(sl, MARGIN, Inches(2.72), Inches(11.5), Inches(1.0),
          "AI Virtual KOL — Status Review", size=42, color=WHITE, bold=True)
    txbox(sl, MARGIN, Inches(3.8), Inches(11.5), Inches(0.6),
          "What works · What is blocking us · What I need decided", size=19,
          color=RGBColor(0xAE, 0xC0, 0xE4))
    txbox(sl, MARGIN, Inches(4.66), Inches(11.5), Inches(0.9),
          f"{today.isoformat()}  ·  every figure measured on this machine  ·  "
          f"reproducible with tools/selftest",
          size=12.5, color=RGBColor(0x8D, 0x9D, 0xBE))

    # 2 — executive summary -------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Executive summary", "The pipeline works end to end. Two quality gaps block launch.")
    stat_row(sl, Inches(1.62), [
        ("4 of 4", "pipeline layers working", OK),
        (f'{M["turn_s"]} s', f'wait per reply — {M["tts_pct"]}% is speech', BAD),
        (f'{M["tuned_pct"]}%', f'persona compliance (was {M["base_pct"]}%)', WARN),
        (f'1 of {s["kols_total"]}', "characters fully built", WARN),
    ])
    bullets(sl, MARGIN, Inches(3.0), BODY_W, [
        ("A viewer can comment and she answers out loud, on camera, in her own voice. ",
         "Her face, her cloned voice, an answer written in her persona — generated locally, "
         "nothing sent to a third party, no per-reply API fee."),
        ("Problem 1 — she still sounds stiff rather than interesting. ",
         "Five separate causes, and the root one is not fixable by more prompt engineering: the "
         "training data is the 7B model's own output, so it cannot teach charm the model does "
         "not already have."),
        ("Problem 2 — she is slow to answer. ",
         f'{M["turn_s"]} seconds per reply, and {M["tts_pct"]}% of that is speech synthesis '
         f'running in a mode that waits for the whole sentence before playing any of it. This '
         f'one is pure engineering and costs nothing to fix.'),
        ("Recommendation: fix the wait first (1–2 weeks, no budget), then the naturalness. ",
         "The naturalness work needs a native English writer and a small one-off budget — "
         "that is the decision I am asking for."),
    ], size=13, gap=Inches(0.7))

    # 3 — architecture ------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "The system", "Four layers, each a separate local service")
    stages = [
        ("1 · Brain", "Writes the reply in persona\nand checks its own output\nagainst hard rules",
         f'{M["llm_s"]} s', OK, "WORKING"),
        ("2 · Voice", "Speaks it in that character's\nown cloned voice",
         f'{M["tts_s"]} s', BAD, "THE BOTTLENECK"),
        ("3 · Face", "Lip-sync, blinking and head\nmotion from one still photo",
         "realtime", WARN, "LICENCE SWAP DUE"),
        ("4 · Broadcast", "WebRTC, OBS or RTMP\nout to the platform",
         "verified", OK, "WORKING"),
    ]
    gap = Inches(0.18)
    cw = int((BODY_W - gap * 3) / 4)
    y = Inches(1.7)
    for i, (name, desc, note, col, badge) in enumerate(stages):
        x = MARGIN + i * (cw + gap)
        panel(sl, x, y, cw, Inches(2.2))
        txbox(sl, x + Inches(0.18), y + Inches(0.16), cw - Inches(0.36), Inches(0.3),
              name, size=14, color=DEEP, bold=True)
        txbox(sl, x + Inches(0.18), y + Inches(0.5), cw - Inches(0.36), Inches(0.26),
              badge, size=9.5, color=col, bold=True)
        txbox(sl, x + Inches(0.18), y + Inches(0.86), cw - Inches(0.36), Inches(0.9),
              desc, size=11, color=MUTED)
        txbox(sl, x + Inches(0.18), y + Inches(1.74), cw - Inches(0.36), Inches(0.3),
              note, size=10.5, color=col, bold=True)
    panel(sl, MARGIN, Inches(4.2), BODY_W, Inches(1.0), fill=WHITE)
    txbox(sl, MARGIN + Inches(0.24), Inches(4.42), BODY_W - Inches(0.5), Inches(0.6),
          "comment  →  brain (local LLM + rule guards)  →  voice  →  lip-sync  →  live video",
          size=13, color=INK, bold=True)
    txbox(sl, MARGIN, Inches(5.42), BODY_W, Inches(1.2),
          "One character profile is the single source of truth: personality, language, voice, "
          "face and selling role all read from it, so adding a character is adding a folder "
          "rather than writing code. The layers talk over HTTP because each engine pins "
          "conflicting libraries — which also means any one of them can be upgraded or replaced "
          "without touching the others.", size=12, color=MUTED)

    # 4 — the comment loop --------------------------------------------------
    #
    # Its own slide because this is the commercial engine: the selling happens in the comment
    # and DM layer, not in the post. It was one line on the evidence slide, which understated
    # the part of the system a manager actually cares about.
    sl = slide_blank(prs)
    header(sl, "Reading comments, answering questions", "The five steps between a comment and her voice")
    steps = [
        ("1 · A comment arrives", "It joins the queue with the viewer's name on it. Two surfaces "
         "use the same brain: the live stream, and a private chat that remembers the person "
         "between visits."),
        ("2 · She picks the register", "Read from the message itself, not from a dropdown — a "
         "song request, someone opening up, or an ordinary comment. On a real stream nobody "
         "labels their comment before sending it."),
        ("3 · Draft, then guards", "The reply is written in persona and rule-checked in code. A "
         "violation is rewritten with the broken rule named; a second failure is never spoken."),
        ("4 · A human approves", "Approve, edit or reject. An edited reply is re-checked — a "
         "person can paste in a price as easily as a model can invent one."),
        ("5 · She says it", "Pushed to the live avatar in her own voice. The next two answers "
         "are already rendering while this one plays."),
    ]
    gap = Inches(0.16)
    cw = int((BODY_W - gap * 4) / 5)
    y = Inches(1.66)
    for i, (title, body) in enumerate(steps):
        x = MARGIN + i * (cw + gap)
        panel(sl, x, y, cw, Inches(1.68))
        txbox(sl, x + Inches(0.16), y + Inches(0.14), cw - Inches(0.32), Inches(0.3),
              title, size=11.5, color=DEEP, bold=True)
        txbox(sl, x + Inches(0.16), y + Inches(0.52), cw - Inches(0.32), Inches(1.05),
              body, size=9.5, color=MUTED)
    rows = [
        ["Capability", "How it works today", "State"],
        ["Two answering surfaces",
         "Live stream: a broadcast queue she works through on her own, holding each answer until "
         "the last has finished playing. Private chat: one person, thread kept, facts about them "
         "remembered between visits", ("Working", OK)],
        ["Approve before send",
         "Every character is set to suggest mode. Draft → approve / edit / reject, recorded in an "
         "append-only log that survives a crash mid-review. A draft the guards refused needs an "
         "edit, not a click", ("Enforced", OK)],
        ["Pulling comments from the platforms",
         "Comments are typed or pasted in today. Instagram Graph and YouTube Live are the "
         "first-party-safe routes, and are not connected yet", ("Not built", WARN)],
    ]
    table(sl, MARGIN, Inches(3.56), BODY_W, rows,
          [Inches(2.7), Inches(7.5), Inches(1.69)], size=10.5, row_h=Inches(0.68))
    txbox(sl, MARGIN, Inches(6.44), BODY_W, Inches(0.8),
          "TikTok is deliberately excluded: it has no official comment or live-chat API, so "
          "automating it would put the account at risk. That is a design decision, not a missing "
          "feature — and it is why the ingestion work is scoped to Instagram and YouTube.",
          size=11.5, color=MUTED, italic=True)

    # 5 — proven ------------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "What is proven", "Verified by measurement, not by impression")
    rows = [
        ["Claim", "Evidence", "Result"],
        ["Five distinct cloned voices",
         f'Speaker verification: cross-character similarity {M["spk_cross"]} against a '
         f'{M["spk_base"]} same-speaker baseline', ("Verified", OK)],
        ["Sofia's voice reaches human range",
         f'Pitch range {M["pitch_range"]} semitones vs {M["pitch_human"]} for the real human '
         f'clip; the old flat voice was {M["pitch_old"]}', ("Verified", OK)],
        ["She says the words we gave her",
         f'Speech transcribed back by a separate recogniser: {M["asr"]}% word accuracy',
         ("Verified", OK)],
        ["Safety rules hold",
         "13-case battery: 9 that must be blocked, 4 that must NOT be — prices, links, "
         "medical claims, jailbreaks, denying being AI", ("0 violations", OK)],
        ["A human approves before anything is sent",
         "Draft → review → approve/edit/reject. A human-edited reply is re-checked too",
         ("Enforced", OK)],
        ["The avatar looks alive",
         f'{M["blinks"]} blinks per minute (natural is 15–20). The procedural alternative '
         f'blinked 0 times', ("Verified", OK)],
        ["It fits the current GPU",
         f'Voice and avatar resident together: {M["vram_used"]:,} MiB used, '
         f'{M["vram_free"]:,} MiB free', ("Verified", OK)],
    ]
    table(sl, MARGIN, Inches(1.62), BODY_W, rows,
          [Inches(3.3), Inches(6.9), Inches(1.69)], size=10.5, row_h=Inches(0.56))
    txbox(sl, MARGIN, Inches(6.28), BODY_W, Inches(0.8),
          "Language coverage was measured the same way rather than assumed: English 1.00, "
          "Japanese 0.80, Spanish 0.64, Vietnamese 0.087. Vietnamese was removed from the "
          "interface — offering it would be offering something broken.",
          size=11.5, color=MUTED, italic=True)

    # 6 — problem A: stiff replies ------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Problem 1", "Her replies are stiff, and not yet interesting")
    txbox(sl, MARGIN, Inches(1.58), BODY_W, Inches(0.4),
          "Five separate causes. Fixing any one of them alone will not move it:",
          size=12.5, color=MUTED)
    items = [
        ("a · The 7B model has a ceiling", BAD,
         "Given every rule explicitly in its prompt, it still greeted a fan again mid-conversation "
         "on 2 turns out of 4 and signed off with an offer of further help on 1 of 4. It writes "
         "the right reply about half the time."),
        ("b · The training data is distilled from itself — the root cause", BAD,
         f'The {M["train_rows"]} training examples are the same 7B model\'s own answers, filtered '
         f'to the ones that obeyed. That fixes habits ({M["tuned_pct"]}% vs {M["base_pct"]}%) but '
         f'cannot add wit the base model does not have. This is the current quality ceiling.'),
        ("c · The safety guards push toward bland", WARN,
         "Break a rule twice and she speaks a safe fallback line — which is the most evasive "
         "sentence in the system. Observed: a viewer asked for a story about her day and was "
         "told \"let me double-check that before I answer\"."),
        ("d · She has no life to talk about", WARN,
         "What separates a person's answer from a model's is a concrete detail. The mechanism "
         "exists (life threads, per-fan memory); what is fed into it is thin."),
        ("e · Nobody has judged it by ear", WARN,
         "Every quality number today comes from an automatic judge built from the same rules "
         "used to filter the training data. It measures rule-breaking, not charm."),
    ]
    y = Inches(2.06)
    for title, col, body in items:
        panel(sl, MARGIN, y, BODY_W, Inches(0.86), fill=WHITE)
        dot = sl.shapes.add_shape(9, MARGIN + Inches(0.18), y + Inches(0.33),
                                  Inches(0.13), Inches(0.13))
        dot.fill.solid(); dot.fill.fore_color.rgb = col
        dot.line.fill.background(); dot.shadow.inherit = False
        txbox(sl, MARGIN + Inches(0.48), y + Inches(0.11), BODY_W - Inches(0.7), Inches(0.28),
              title, size=12.5, color=DEEP, bold=True)
        txbox(sl, MARGIN + Inches(0.48), y + Inches(0.42), BODY_W - Inches(0.75), Inches(0.4),
              body, size=10.5, color=MUTED)
        y += Inches(0.95)

    # 7 — problem B: latency ------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Problem 2", f'She takes {M["turn_s"]} seconds to answer — and it is not the thinking')
    # The split, drawn to scale. The point of the slide is the proportion, so the bar is the
    # argument and the numbers only label it.
    bw = BODY_W
    y = Inches(1.72)
    llm_w = int(bw * M["llm_s"] / M["turn_s"])
    b1 = sl.shapes.add_shape(1, MARGIN, y, llm_w, Inches(0.62))
    b1.fill.solid(); b1.fill.fore_color.rgb = ACCENT
    b1.line.fill.background(); b1.shadow.inherit = False
    b2 = sl.shapes.add_shape(1, MARGIN + llm_w, y, bw - llm_w, Inches(0.62))
    b2.fill.solid(); b2.fill.fore_color.rgb = BAD
    b2.line.fill.background(); b2.shadow.inherit = False
    txbox(sl, MARGIN + llm_w + Inches(0.18), y + Inches(0.16), Inches(6), Inches(0.34),
          f'{M["tts_s"]} s  ·  speech synthesis  ·  {M["tts_pct"]}% of the wait',
          size=14, color=WHITE, bold=True)
    txbox(sl, MARGIN, y + Inches(0.74), bw, Inches(0.3),
          f'▲ thinking: {M["llm_s"]} s', size=11, color=ACCENT, bold=True)
    rows = [
        ["Why it is slow", "Detail"],
        ["Speech runs in non-streaming mode",
         "The whole sentence must finish before the first sound plays — even though the engine "
         "produces audio FASTER than it can be listened to (real-time factor 0.54–0.68)"],
        ["Private chat cannot pre-render",
         "The live stream already renders the next two answers while the current one plays. A "
         "one-to-one chat does not know the question in advance"],
        ["Each retry is another model call",
         "A rule violation or a generic answer triggers a rewrite; worst case is three calls "
         "for one reply"],
    ]
    table(sl, MARGIN, Inches(3.06), BODY_W, rows,
          [Inches(3.9), Inches(7.99)], size=11, row_h=Inches(0.62))
    txbox(sl, MARGIN, Inches(5.72), BODY_W, Inches(1.2),
          "The good news is in the second column: the engine is already fast enough. It is being "
          "asked to work in the wrong mode. Streaming the speech sentence by sentence — playing "
          "the first while the second is still being made — is the single highest-value fix in "
          "this deck, and it needs no budget and no new hardware.",
          size=12.5, color=INK)

    # 8 — problem C: system, assets, licence --------------------------------
    sl = slide_blank(prs)
    header(sl, "Problem 3", "Hardware, licences and unfinished characters")
    rows = [
        ["Issue", "Detail", "Blocks"],
        ["12 GB VRAM is a hard ceiling",
         f'{gpu.get("name", "RTX 5070")} — enough to run, not enough to run everything at once '
         f'or to train a lip-sync model (needs 23–30 GB)', ("Scale", WARN)],
        ["wav2lip weights are research-licensed",
         "The lip-sync model in use today cannot be used commercially. MuseTalk was licence-"
         "checked and is clean all the way down (MIT, weights explicitly commercial) and our "
         "stack already supports it — a flag change plus half a day of testing",
         ("Commercial launch", BAD)],
        ["The lip-sync server asks for its watermark",
         "0.8% of the frame, low-contrast grey. Accepted for now; one email to the author would "
         "replace the ambiguity with an answer", ("Publishing", WARN)],
        ["Only one character is complete",
         f'{s["kols_with_images"]} of {s["kols_total"]} have any images at all; '
         f'{s["kols_voiced"]} have a finished voice. The flagship seller\'s product catalogue is '
         f'still an empty template', ("Roster growth", WARN)],
    ]
    table(sl, MARGIN, Inches(1.66), BODY_W, rows,
          [Inches(3.2), Inches(6.9), Inches(1.79)], size=10.5, row_h=Inches(0.74))
    txbox(sl, MARGIN, Inches(5.62), BODY_W, Inches(1.5),
          "Six components in this layer have now been licence-checked before adoption and four "
          "carried a restriction — so this is treated as a standing checklist item rather than a "
          "surprise. None of it blocks development; all of it blocks publishing, which is why "
          "the swap is scheduled before the first commercial output rather than after.",
          size=12, color=MUTED)

    # 9 — fixes P0 ----------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Fix 1 — the wait", "Engineering only. No budget, no new hardware, 1–2 weeks.")
    rows = [
        ["#", "Action", "Expected effect", "Cost"],
        ["1", "Stream the speech sentence by sentence — play the first while the second renders",
         ("Time to first sound: 7.6 s → ~1.5–2 s", OK), "1–2 weeks"],
        ["2", "Feed the model's output into speech as it is written, at the first full stop",
         "Removes most of the remaining thinking time", "same work"],
        ["3", "Merge the two retry passes into one, checking the draft as it streams",
         "Removes the worst-case triple call", "a few days"],
        ["4", "Enable fp16 / JIT / TensorRT on the speech engine (all three are off today)",
         ("Needs measuring — possibly 1.3–2×", WARN), "half a day"],
        ["5", "Size the reply to the question — a six-sentence answer is ~10 s of audio",
         "Less waiting, better pacing", "a day"],
    ]
    table(sl, MARGIN, Inches(1.66), BODY_W, rows,
          [Inches(0.5), Inches(5.6), Inches(4.2), Inches(1.59)], size=11, row_h=Inches(0.64))
    txbox(sl, MARGIN, Inches(5.62), BODY_W, Inches(1.4),
          "Items 1–3 are the certain part; 4 and 5 need measuring before I would promise them. "
          "The target is a reply that starts speaking in under two seconds, which is the "
          "threshold where a conversation stops feeling like a machine taking its turn.\n\n"
          "Nothing on this slide needs a decision — it is listed so the plan is visible.",
          size=12, color=MUTED)

    # 10 — fixes P1: the four levers ------------------------------------------
    #
    # Rewritten from an action list into the mechanism. The question this slide is actually
    # asked is "HOW does she get more interesting", and a numbered to-do list does not answer
    # it — it assumes the answer. Ordered by value per unit of cost, not by sequence.
    sl = slide_blank(prs)
    header(sl, "Fix 2 — how she learns to sound like a person",
           "Four levers. The cheapest is also the strongest, and it is not the model.")
    rows = [
        ["", "Lever", "What it actually does", "Cost"],
        ["1", "Give her something to say",
         "The gap between a person and a model is a FACT, not a capability. \"That sounds hard\" "
         "is a model; \"I had a week like that in March and lived on toast\" is a person. A "
         "weekly journal of real specifics — what she filmed, what went wrong — plus what she "
         "already remembers about each fan. A 7B model can be specific if it is given specifics",
         ("1–2 h/week", OK)],
        ["2", "Show, don't tell",
         "Measured here repeatedly: every rule added changed the wording and left the behaviour. "
         "Ban \"take a walk\" and the next reply is \"maybe a change of scenery\" — the same "
         "shrug in new clothes. Concrete example replies changed it. 200–300 written by a native "
         "speaker define what interesting means for THIS character",
         ("A writer", BAD)],
        ["3", "Raise the ceiling of the training data",
         "The pipeline already generates candidate replies and filters them with the production "
         "guards. Today the generator is the same 7B being trained, so its ceiling is its own "
         "best half. Point the generator at a stronger model and keep the same judge — existing "
         "code, one changed setting", ("Small one-off API budget", WARN)],
        ["4", "Teach the axis rules cannot describe",
         "A rule can say \"no bullet points\". Only a preference can say \"this one is funnier\". "
         "Keep the better/worse pairs from the blind scoring and train on the preference itself",
         ("Needs 1–2 done first", MUTED)],
    ]
    table(sl, MARGIN, Inches(1.62), BODY_W, rows,
          [Inches(0.4), Inches(2.5), Inches(7.2), Inches(1.79)], size=10, row_h=Inches(0.92))
    txbox(sl, MARGIN, Inches(6.36), BODY_W, Inches(0.75),
          "None of it is provable without scoring by ear first — 20 questions, 3 people, half a "
          "day, before and after. Levers 1 and 2 are where I would start: they need no new "
          "model, no new hardware, and they are what makes her specific rather than pleasant.",
          size=12, color=INK, bold=True)

    # 11 — the two obvious shortcuts -----------------------------------------
    #
    # Both of these get proposed in every conversation about this, so they are answered on a
    # slide rather than in the room. Neither is wrong; both are second.
    sl = slide_blank(prs)
    header(sl, "The two obvious shortcuts", "Both get suggested every time. Here is what each is really worth.")

    panel(sl, MARGIN, Inches(1.62), BODY_W, Inches(2.62), fill=WHITE)
    txbox(sl, MARGIN + Inches(0.26), Inches(1.76), BODY_W - Inches(0.5), Inches(0.3),
          "“Train her on real streamers and YouTubers, let her learn by herself”",
          size=14, color=DEEP, bold=True)
    txbox(sl, MARGIN + Inches(0.26), Inches(2.12), BODY_W - Inches(0.55), Inches(0.3),
          "What it would genuinely give us is the right thing: the shape of real spoken "
          "conversation — how long a reply runs, how often a real person names a detail or takes "
          "a side. That is exactly what is missing. Three problems stop it being training data:",
          size=10.5, color=MUTED)
    for i, (mk, line) in enumerate([
        ("a", "A transcript is a monologue, not question-and-answer pairs. Training needs "
              "\"viewer said X → she replied Y\", which means aligning the live chat to the "
              "speech timeline — and that replay does not exist on every platform."),
        ("b", "Rights. Training a commercial persona on a named creator's words is the same "
              "category of exposure as cloning their voice — which this project already refuses "
              "to do without consent. A legal question, not a technical one."),
        ("c", "It teaches her to sound like THEM. She would inherit somebody else's persona, "
              "which is the opposite of the goal."),
    ]):
        yy = Inches(2.66) + Inches(0.34) * i
        txbox(sl, MARGIN + Inches(0.3), yy, Inches(0.3), Inches(0.26), mk,
              size=10.5, color=ACCENT, bold=True)
        txbox(sl, MARGIN + Inches(0.6), yy, BODY_W - Inches(0.9), Inches(0.3), line,
              size=10.5, color=MUTED)
    txbox(sl, MARGIN + Inches(0.26), Inches(3.78), BODY_W - Inches(0.55), Inches(0.34),
          "The version that does work: mine them for the RUBRIC, not for the answers. Measure "
          "real streams — reply length, how often they ask something back, how often a concrete "
          "detail appears — and use those numbers to score our own output and brief the writer.",
          size=11, color=ACCENT, bold=True)

    panel(sl, MARGIN, Inches(4.42), BODY_W, Inches(1.72), fill=WHITE)
    txbox(sl, MARGIN + Inches(0.26), Inches(4.56), BODY_W - Inches(0.5), Inches(0.3),
          "“Just use a bigger model”", size=14, color=DEEP, bold=True)
    txbox(sl, MARGIN + Inches(0.26), Inches(4.92), BODY_W - Inches(0.55), Inches(0.6),
          "Real, but second — and it is gated by the GPU rather than by the idea. The 7B in use "
          "takes about 4.7 GB, which fits in the 6.6 GB left after the voice and the avatar are "
          "loaded. A 14B at the same quantisation is roughly 9 GB: it does not fit alongside "
          "them on a 12 GB card, so it means either a 24 GB card or running the layers in turns. "
          "It is also slower per reply, which pulls against the wait we are trying to remove.",
          size=10.5, color=MUTED)
    txbox(sl, MARGIN + Inches(0.26), Inches(5.66), BODY_W - Inches(0.55), Inches(0.34),
          "Verdict: test it the week the GPU question is answered — and do not wait for it "
          "before starting on levers 1 and 2.", size=11, color=ACCENT, bold=True)

    txbox(sl, MARGIN, Inches(6.3), BODY_W, Inches(0.9),
          "Both shortcuts share the same flaw: they answer \"where do we get more data\" when "
          "the binding constraint is \"what should the data be an example OF\". Nobody has "
          "written down what an interesting reply from this character looks like — that is the "
          "work, and it is why the ask is a writer rather than a scraper.",
          size=11.5, color=MUTED, italic=True)

    # 12 — SUPPORT NEEDED, numbered -----------------------------------------
    sl = slide_blank(prs)
    header(sl, "What I need from you", "Three decisions. Everything else I can do myself.")
    asks = [
        ("1", "A part-time native English writer", BAD,
         "To write 200–300 sample replies and to help score the blind test. This is the single "
         "biggest lever on \"she sounds stiff\" — and the only item on this list that no amount "
         "of engineering can substitute for.",
         "Roughly 2–3 days of their time, spread over two weeks."),
        ("2", "A small one-off budget for teacher-model API calls", WARN,
         "Used once, to generate high-quality training data that the local 7B model cannot "
         "produce for itself. Small in absolute terms, but it needs approval to spend.",
         "One-off. I will report exactly what was spent and what it changed."),
        ("3", "A decision on the 24 GB GPU — buy or rent?", WARN,
         "Unlocks the 14B model, lets all three layers run at once instead of taking turns, and "
         "makes training our own lip-sync model possible. The cost model is already in the "
         "executive deck.",
         "Not blocking today's work; it is the ceiling on everything after it."),
    ]
    y = Inches(1.66)
    for num, title, col, body, note in asks:
        panel(sl, MARGIN, y, BODY_W, Inches(1.42), fill=WHITE)
        badge = sl.shapes.add_shape(9, MARGIN + Inches(0.24), y + Inches(0.3),
                                    Inches(0.62), Inches(0.62))
        badge.fill.solid(); badge.fill.fore_color.rgb = col
        badge.line.fill.background(); badge.shadow.inherit = False
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.size = Pt(24); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
        txbox(sl, MARGIN + Inches(1.12), y + Inches(0.18), BODY_W - Inches(1.4), Inches(0.32),
              title, size=15, color=DEEP, bold=True)
        txbox(sl, MARGIN + Inches(1.12), y + Inches(0.56), BODY_W - Inches(1.45), Inches(0.5),
              body, size=11, color=MUTED)
        txbox(sl, MARGIN + Inches(1.12), y + Inches(1.08), BODY_W - Inches(1.45), Inches(0.26),
              note, size=10.5, color=col, bold=True)
        y += Inches(1.55)
    txbox(sl, MARGIN, Inches(6.42), BODY_W, Inches(0.6),
          "Two more that are not decisions but need someone else's eyes: legal to read the "
          "MuseTalk licence (so we can drop the research-licensed model), and a call on how "
          "intimate the companion mode is allowed to be — that is a brand question, not a "
          "technical one.", size=11.5, color=MUTED, italic=True)

    # 13 — roster -----------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "The roster", "One character deep, nine wide")
    rows = [["Character", "Images", "Video", "Voice", "State"]]
    order = {"active": 0, "photos_collected": 1, "draft": 2}
    for k in sorted(st["kols"], key=lambda k: (order.get(k["status"], 3), -k["images"])):
        v = k["voice"]
        voice = ("fine-tuned" if v.get("sovits") and v.get("gpt") else
                 "dataset ready" if v.get("clips") else "—")
        if k["id"] == "sofia-vargas":
            voice = "CosyVoice 2"
        state_cell = ("Complete — reference build", OK) if k["id"] == "sofia-vargas" else \
                     ("Images ready", OK) if k["images"] >= 20 else \
                     ("Voice ready, no images", WARN) if voice == "fine-tuned" else \
                     ("Started", WARN) if k["images"] else ("Profile only", MUTED)
        rows.append([k["name"], str(k["images"]), str(k["videos"]), voice, state_cell])
    table(sl, MARGIN, Inches(1.62), BODY_W, rows,
          [Inches(3.4), Inches(1.1), Inches(1.1), Inches(2.4), Inches(3.89)],
          size=10.5, row_h=Inches(0.38))
    txbox(sl, MARGIN, Inches(1.62) + Inches(0.38) * len(rows) + Inches(0.22), BODY_W, Inches(1.2),
          "My recommendation: go deep on Sofia until one character is genuinely good enough to "
          "publish, rather than wide. Every improvement to the brain and the voice is reused by "
          "the next character — the second one cost about 25 minutes of compute and half an hour "
          "of review, against days for the first — whereas adding characters now would replicate "
          "exactly the weaknesses described earlier.", size=12, color=INK)

    # 14 — roadmap ----------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Six-week plan", "Each week ends in something measurable")
    rows = [
        ["Week", "Focus", "How we will know it worked"],
        ["1–2", "Stream the speech; merge the retry passes",
         ("First sound in under 2 seconds, measured", OK)],
        ["2", "Blind A/B test scored by people",
         ("A human baseline for \"does she sound natural\"", OK)],
        ["3–4", "Native-written replies + teacher-model data + the life journal",
         "300+ reviewed training examples"],
        ["5", "Retrain, then preference-train on the blind-test pairs",
         "Second blind score, compared against week 2"],
        ["6", "Swap in MuseTalk; record the full demo",
         ("A demo that is clear to publish commercially", OK)],
    ]
    table(sl, MARGIN, Inches(1.7), BODY_W, rows,
          [Inches(1.2), Inches(5.6), Inches(5.09)], size=11.5, row_h=Inches(0.62))
    txbox(sl, MARGIN, Inches(5.62), BODY_W, Inches(1.4),
          "The two workstreams are independent: the speed work is mine alone and starts "
          "immediately, while the personality work only starts once the writer is available. "
          "If the answer on item 1 is no, weeks 3–5 change shape — I would fall back to "
          "teacher-model data alone, which is cheaper and measurably weaker.",
          size=12, color=MUTED)

    # 15 — close ------------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "In one page", "Where this stands")
    stat_row(sl, Inches(1.66), [
        ("Works", "all four layers, verified", OK),
        (f'{M["turn_s"]} s → 2 s', "the speed fix, no budget", ACCENT),
        ("Needs a writer", "the naturalness fix", BAD),
        ("NT$0", "software licence cost", OK),
    ])
    bullets(sl, MARGIN, Inches(3.1), BODY_W, [
        ("The hard engineering is done and it is measured, not claimed. ",
         "Voice, safety, avatar and the local brain all have evidence behind them, and the "
         "whole thing re-verifies itself with one command."),
        ("What is left is a quality problem, and quality needs a human in the loop. ",
         "A model can be taught to stop breaking rules by another model. It cannot be taught to "
         "be interesting by one."),
        ("The speed fix starts this week regardless of any decision. ",
         "It is the one that changes how the product feels, and it is free."),
    ], size=13, gap=Inches(0.72))
    panel(sl, MARGIN, Inches(5.5), BODY_W, Inches(1.1), fill=WHITE)
    txbox(sl, MARGIN + Inches(0.3), Inches(5.72), BODY_W - Inches(0.6), Inches(0.7),
          "Asking for: 1 · a part-time native English writer   ·   2 · a small one-off API "
          "budget   ·   3 · a decision on the 24 GB GPU",
          size=14, color=DEEP, bold=True, align=PP_ALIGN.CENTER)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("-o", "--out", default=str(REPO / "docs" / "AI-KOL-Status-Review.pptx"))
    args = ap.parse_args()

    import server  # tools/dashboard/server.py
    state = server.collect()
    out = build(state, Path(args.out))
    print(f"wrote {out}  ({out.stat().st_size/1e3:.0f} KB)")
    print(f"15 slides · {M['turn_s']} s per turn ({M['tts_pct']}% speech) · "
          f"{state['summary']['kols_total']} characters")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
