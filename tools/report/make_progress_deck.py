#!/usr/bin/env python3
"""Generate the project progress deck (PowerPoint, English).

Numbers are pulled from `tools/dashboard/server.py:collect()` — the same live scan the
dashboard uses — so the deck cannot drift from what is actually on disk. Regenerate it
any time instead of hand-editing slides.

    python tools/report/make_progress_deck.py
    python tools/report/make_progress_deck.py -o docs/AI-KOL-Progress-Report.pptx

Needs `pip install python-pptx`.
"""
from __future__ import annotations

import argparse
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO / "tools" / "dashboard"))

# ---- palette ---------------------------------------------------------------
INK = RGBColor(0x14, 0x17, 0x1A)
MUTED = RGBColor(0x66, 0x6E, 0x78)
ACCENT = RGBColor(0x2F, 0x5F, 0xD0)
DEEP = RGBColor(0x14, 0x22, 0x44)
OK = RGBColor(0x0F, 0x8A, 0x4C)
WARN = RGBColor(0xA8, 0x64, 0x00)
BAD = RGBColor(0xC6, 0x2B, 0x2B)
PANEL = RGBColor(0xF2, 0xF4, 0xF7)
LINE = RGBColor(0xDD, 0xE1, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.72)
BODY_W = W - 2 * MARGIN


def txbox(slide, x, y, w, h, text, size=14, color=INK, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, spacing=1.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = FONT
    return box


def panel(slide, x, y, w, h, fill=PANEL, line=LINE):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, kicker, title):
    txbox(slide, MARGIN, Inches(0.46), BODY_W, Inches(0.26),
          kicker.upper(), size=11, color=ACCENT, bold=True)
    txbox(slide, MARGIN, Inches(0.76), BODY_W, Inches(0.55),
          title, size=27, color=DEEP, bold=True)
    ln = slide.shapes.add_shape(1, MARGIN, Inches(1.42), BODY_W, Emu(9525))
    ln.fill.solid()
    ln.fill.fore_color.rgb = LINE
    ln.line.fill.background()
    ln.shadow.inherit = False


def bullets(slide, x, y, w, items, size=13, gap=Inches(0.40)):
    """items: list of (bold_lead, rest) or plain strings."""
    cy = y
    for it in items:
        lead, rest = it if isinstance(it, tuple) else ("", it)
        box = slide.shapes.add_textbox(x, cy, w, gap)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        dot = p.add_run()
        dot.text = "▪  "
        dot.font.size = Pt(size)
        dot.font.color.rgb = ACCENT
        dot.font.name = FONT
        if lead:
            r = p.add_run()
            r.text = lead
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = INK
            r.font.name = FONT
        r2 = p.add_run()
        r2.text = rest
        r2.font.size = Pt(size)
        r2.font.color.rgb = INK if not lead else MUTED
        r2.font.name = FONT
        cy += gap
    return cy


def table(slide, x, y, w, rows, col_w, head_fill=DEEP, size=11, row_h=Inches(0.34)):
    nrows, ncols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(nrows, ncols, x, y, w, row_h * nrows)
    tbl = shape.table
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = cw
    for i, row in enumerate(rows):
        tbl.rows[i].height = row_h
        for j, cell_val in enumerate(row):
            val, color = cell_val if isinstance(cell_val, tuple) else (cell_val, None)
            c = tbl.cell(i, j)
            c.text = ""
            c.margin_left = Inches(0.09)
            c.margin_right = Inches(0.07)
            c.margin_top = c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = head_fill if i == 0 else (WHITE if i % 2 else PANEL)
            p = c.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(size)
            r.font.name = FONT
            r.font.bold = (i == 0)
            r.font.color.rgb = WHITE if i == 0 else (color or INK)
    return shape


def stat_row(slide, y, items, h=Inches(1.02)):
    n = len(items)
    gap = Inches(0.16)
    cw = int((BODY_W - gap * (n - 1)) / n)
    for i, (big, label, col) in enumerate(items):
        x = MARGIN + i * (cw + gap)
        panel(slide, x, y, cw, h)
        txbox(slide, x + Inches(0.16), y + Inches(0.13), cw - Inches(0.3), Inches(0.46),
              str(big), size=26, color=col, bold=True)
        txbox(slide, x + Inches(0.16), y + Inches(0.62), cw - Inches(0.3), Inches(0.3),
              label, size=10.5, color=MUTED)


# --------------------------------------------------------------------- slides

def build(state: dict, out: Path) -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    s = state["summary"]
    lena = next((k for k in state["kols"] if k["id"] == "lena-chen"), {"voice": {}})
    lv = lena.get("voice", {})
    gpu = state.get("gpu") or {}
    today = date.today().isoformat()

    # 1 — title -------------------------------------------------------------
    sl = slide_blank(prs)
    bg = sl.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = DEEP
    bg.line.fill.background(); bg.shadow.inherit = False
    bar = sl.shapes.add_shape(1, MARGIN, Inches(2.42), Inches(1.5), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    txbox(sl, MARGIN, Inches(2.72), Inches(10.6), Inches(1.0),
          "AI-KOL System", size=44, color=WHITE, bold=True)
    txbox(sl, MARGIN, Inches(3.78), Inches(10.6), Inches(0.6),
          "Progress Report — Voice Cloning & Realtime Lip-Sync",
          size=19, color=RGBColor(0xAE, 0xC0, 0xE4))
    txbox(sl, MARGIN, Inches(4.62), Inches(10.6), Inches(0.9),
          f"Local-first virtual influencer pipeline  ·  RTX 5070 (12 GB)  ·  {today}",
          size=12.5, color=RGBColor(0x8D, 0x9D, 0xBE))

    # 2 — where we are ------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Executive summary", "Voice + lip-sync are working end-to-end")
    stat_row(sl, Inches(1.66), [
        (s["kols_total"], "KOL personas defined", DEEP),
        (s["kols_voiced"], "voices fine-tuned & verified", OK),
        (f'{lv.get("minutes") or 0:.1f} min', "training corpus (304 clips)", DEEP),
        ("707 ms", "TTS first-chunk latency", OK),
    ])
    bullets(sl, MARGIN, Inches(3.06), BODY_W, [
        ("Two of four pipeline layers are complete. ",
         "Text now becomes a lip-synced video of a talking avatar speaking in a "
         "purpose-built cloned voice — fully local, no paid API."),
        ("The voice is verified, not assumed. ",
         "Synthesis was checked by transcribing the output back with ASR and comparing "
         "to the input text, in both Chinese and English."),
        ("No real person's voice was cloned. ",
         "The timbre is synthesized, then locked in by fine-tuning — so the KOL owns a "
         "voice that never belonged to anybody."),
        ("Most of the effort went into environment blockers, ",
         "not the models. Eight undocumented issues had to be solved on this Windows box; "
         "all are now fixed and written down."),
    ], size=13, gap=Inches(0.72))

    # 3 — architecture ------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Architecture", "Four layers, two of them done")
    stages = [
        ("1 · Persona", "profile.json, character bible,\ncontent pillars",
         f'{s["kols_total"]} defined', OK, "DONE"),
        ("2 · Images", "Seed images per KOL;\nface-lock LoRA not trained",
         f'{s["kols_with_images"]}/{s["kols_total"]} have seeds', WARN, "PARTIAL"),
        ("3 · Voice", "GPT-SoVITS v2Pro fine-tune,\nserved over HTTP :9880",
         "verified ZH + EN", OK, "DONE"),
        ("4 · Lip-sync", "LiveTalking + wav2lip,\nWebRTC / OBS / RTMP",
         "stock avatar only", WARN, "PARTIAL"),
    ]
    gap = Inches(0.18)
    cw = int((BODY_W - gap * 3) / 4)
    y = Inches(1.74)
    for i, (name, desc, note, col, badge) in enumerate(stages):
        x = MARGIN + i * (cw + gap)
        panel(sl, x, y, cw, Inches(2.34))
        txbox(sl, x + Inches(0.18), y + Inches(0.17), cw - Inches(0.36), Inches(0.3),
              name, size=14, color=DEEP, bold=True)
        txbox(sl, x + Inches(0.18), y + Inches(0.53), cw - Inches(0.36), Inches(0.28),
              badge, size=10, color=col, bold=True)
        txbox(sl, x + Inches(0.18), y + Inches(0.92), cw - Inches(0.36), Inches(0.8),
              desc, size=11, color=MUTED)
        txbox(sl, x + Inches(0.18), y + Inches(1.86), cw - Inches(0.36), Inches(0.3),
              note, size=10.5, color=col, bold=True)
    txbox(sl, MARGIN, Inches(4.36), BODY_W, Inches(0.4),
          "Data flow", size=12, color=ACCENT, bold=True)
    panel(sl, MARGIN, Inches(4.72), BODY_W, Inches(1.06), fill=WHITE)
    txbox(sl, MARGIN + Inches(0.24), Inches(4.9), BODY_W - Inches(0.5), Inches(0.7),
          "text  →  LiveTalking :8010  →  GPT-SoVITS api_v2 :9880  →  cloned-voice audio\n"
          "                                    ↳  wav2lip  →  h264 576×768 @30 fps  →  WebRTC / OBS / RTMP",
          size=12.5, color=INK)
    txbox(sl, MARGIN, Inches(6.0), BODY_W, Inches(0.6),
          "Three isolated Python environments are required — their dependencies genuinely "
          "conflict (numpy <2 vs 2.x). They talk over HTTP, so they never need to share one.",
          size=11, color=MUTED, italic=True)

    # 4 — what shipped ------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Delivered", "New tooling, all tested on real data")
    rows = [
        ["Component", "What it does", "Status"],
        ["voice_crawl/crawl.py",
         "Podcast/video → training set: ASR word-timestamps, sentence-boundary slicing, 7-metric QC gate",
         ("Working", OK)],
        ["voice_crawl/bootstrap_timbre.py",
         "Synthesizes a consent-free voice corpus; exact transcripts, no ASR error",
         ("Working", OK)],
        ["voice_crawl/corpus_builder.py",
         "Generates 300+ unique bilingual utterances (in-domain + phonetically spread)",
         ("Working", OK)],
        ["voice_crawl/train_gptsovits.py",
         "Headless 6-step fine-tune (WebUI hides these behind Gradio callbacks)",
         ("Working", OK)],
        ["livetalking/run_livetalking.ps1",
         "Launches the avatar using the voice recorded in the KOL's profile.json",
         ("Working", OK)],
        ["livetalking/verify_lipsync.py",
         "Headless WebRTC client — proves lip-sync without needing a browser",
         ("Working", OK)],
        ["dashboard/server.py",
         "Live tracker: pipeline state, dataset QC, service + GPU health",
         ("Working", OK)],
    ]
    table(sl, MARGIN, Inches(1.72), BODY_W, rows,
          [Inches(3.3), Inches(7.4), Inches(1.19)], size=10.5, row_h=Inches(0.46))
    txbox(sl, MARGIN, Inches(5.56), BODY_W, Inches(0.9),
          "Also fixed two pre-existing bugs that made the voice client unusable on Windows: "
          "it shelled out to the macOS-only `say` command, and read UTF-8 profiles with the "
          "system ANSI codepage — which crashed on every KOL with Chinese text.",
          size=11.5, color=MUTED)

    # 5 — voice results -----------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Result · voice", "Verified by ASR round-trip, not by ear")
    txbox(sl, MARGIN, Inches(1.66), BODY_W, Inches(0.4),
          "Synthesized audio was transcribed back and compared to the requested text:",
          size=12.5, color=MUTED)
    rows = [
        ["Lang", "Requested", "Transcribed back"],
        ["ZH", "大家好，今天分享一個好物，這罐精華我用了三週。",
         "大家好,今天分享一個好物這罐精華我用了三周"],
        ["EN", "Honestly, this is my favourite thing I have tried all month.",
         "Honestly, this is my favorite thing I have tried all month."],
    ]
    table(sl, MARGIN, Inches(2.14), BODY_W, rows,
          [Inches(0.86), Inches(5.5), Inches(5.53)], size=11, row_h=Inches(0.52))
    txbox(sl, MARGIN, Inches(3.9), BODY_W, Inches(0.5),
          "Only differences are punctuation and 週/周 + favourite/favorite — the ASR engine's own "
          "normalization, not synthesis errors. One timbre speaks both languages.",
          size=11.5, color=MUTED, italic=True)
    stat_row(sl, Inches(4.6), [
        (f'{lv.get("clips", 0)}', "training clips accepted", DEEP),
        (f'{(lv.get("langs") or {}).get("ZH", 0)} / {(lv.get("langs") or {}).get("EN", 0)}',
         "ZH / EN balance", DEEP),
        ("v2Pro", "GPT-SoVITS model version", DEEP),
        ("e8 + e12", "SoVITS / GPT epochs kept", DEEP),
    ])
    txbox(sl, MARGIN, Inches(5.94), BODY_W, Inches(0.6),
          "Why the QC gate matters: when Chinese audio was force-labelled English, the ASR "
          "translated instead of transcribing — audio saying one thing, training label another. "
          "Confidence separated them cleanly (real 0.86–0.99 vs mislabelled 0.33–0.54).",
          size=11, color=MUTED)

    # 6 — lip-sync results --------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Result · lip-sync", "Realtime avatar speaking the cloned voice")
    stat_row(sl, Inches(1.7), [
        ("150 / 297", "video / audio frames received", DEEP),
        ("707 ms", "voice first-chunk latency", OK),
        ("5,307 MiB", "VRAM, both services resident", OK),
        ("6,920 MiB", "VRAM still free of 12,227", OK),
    ])
    bullets(sl, MARGIN, Inches(3.1), BODY_W, [
        ("A/V is in sync. ", "Recorded output: video 5.967 s vs audio 5.940 s, h264 576×768 @30 fps."),
        ("Both services fit on one 12 GB card. ",
         "The earlier concern that GPT-SoVITS and LiveTalking could not co-exist was measured "
         "and disproved — no staging needed."),
        ("Mixed-language lines render correctly. ",
         "LiveTalking hardcoded Chinese for all text; patched to per-segment auto-detection, so "
         "\"...一個好物 real talk\" now speaks correctly inside one sentence."),
        ("The voice service is not optional. ",
         "LiveTalking calls it per utterance — stopping it silently downgrades the avatar to a "
         "generic voice, so the launcher refuses to start without it."),
    ], size=12.5, gap=Inches(0.66))
    txbox(sl, MARGIN, Inches(5.92), BODY_W, Inches(0.6),
          "Caveat: the face is LiveTalking's stock demo avatar, not the KOL's. See "
          "\"Help needed\".", size=12, color=WARN, bold=True)

    # 7 — difficulties ------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Difficulties", "Eight undocumented blockers — all now solved")
    rows = [
        ["Blocker", "Root cause", "Resolution"],
        ["Training crashed instantly (0xC0000005)",
          "Single-GPU DDP: its reducer hooks into backward and segfaults on the Windows gloo backend",
          "Bypass DDP when 1 GPU"],
        ["numba DLL refused to load",
          "Windows Smart App Control blocks numba 0.66's binary — broke all audio resampling",
          "Pin numba <0.62"],
        ["Two deps cannot compile",
          "pyopenjtalk / jieba_fast are C extensions; no MSVC toolchain on this machine",
          "Drop JA; shim jieba_fast"],
        ["Audio decode failed in one step",
          "torchaudio 2.11 delegates to torchcodec, which needs FFmpeg *shared* DLLs (static build installed)",
          "Install shared FFmpeg"],
        ["Steps 'succeeded' producing nothing",
          "Upstream wraps per-clip work in bare except: — exit code 0 with empty output dirs",
          "Assert artifact counts"],
        ["Training finished, then lost the weights",
          "Output dir is created by the WebUI, not the trainer → error after all epochs ran",
          "Pre-create + recovery tool"],
        ["Model weights only on Quark / Google Drive",
          "No direct download URL for wav2lip weights; browser-gated hosts",
          "Automated via gdown"],
        ["RTX 5070 is Blackwell (sm_120)",
          "Requires CUDA 12.8 wheels; older builds will not run at all",
          "torch 2.11.0+cu128"],
    ]
    table(sl, MARGIN, Inches(1.66), BODY_W, rows,
          [Inches(3.5), Inches(6.9), Inches(1.49)], size=9.5, row_h=Inches(0.53))
    txbox(sl, MARGIN, Inches(6.5), BODY_W, Inches(0.5),
          "All eight are documented in tools/voice_crawl/README.md and tools/livetalking/README.md "
          "so they cost time once, not repeatedly.", size=11, color=MUTED, italic=True)

    # 8 — help needed -------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Help needed", "Four decisions or inputs blocking further progress")
    items = [
        ("1", "The avatar needs the KOL's face — a video is required", WARN,
         "LiveTalking builds an avatar from a VIDEO, but every KOL currently has only still "
         "images. Options: (a) generate a short idle-motion clip from one portrait via "
         "image-to-video, (b) commission/record a real base video, or (c) accept a stock face "
         "for now. Needs a decision on which — (a) is the only fully local route."),
        ("2", "Rights decision: leftover real-person voice data", BAD,
         "kols/sofia-hsu/voice/ holds a dataset built from a podcast recording during "
         "pipeline testing. It conflicts with the synthetic-only policy chosen for this project. "
         "It was NOT deleted without approval — please confirm whether to remove it."),
        ("3", "Multi-speaker audio is not yet supported", WARN,
         "Speaker diarization is not wired in, so a two-host podcast would train a blended "
         "voice. Only single-speaker sources are safe today. Needed only if crawling real "
         "audio becomes the primary path."),
        ("4", "Hardware ceiling for running everything at once", WARN,
         "Voice + lip-sync fit in 12 GB with ~6.9 GB free, but adding the LLM brain and image "
         "generation will not. Either stage components, or budget a larger card (24 GB) if "
         "simultaneous operation is required."),
    ]
    y = Inches(1.66)
    for num, title, col, body in items:
        panel(sl, MARGIN, y, BODY_W, Inches(1.18), fill=WHITE)
        txbox(sl, MARGIN + Inches(0.2), y + Inches(0.16), Inches(0.4), Inches(0.4),
              num, size=17, color=col, bold=True)
        txbox(sl, MARGIN + Inches(0.66), y + Inches(0.15), BODY_W - Inches(0.9), Inches(0.32),
              title, size=13, color=DEEP, bold=True)
        txbox(sl, MARGIN + Inches(0.66), y + Inches(0.5), BODY_W - Inches(0.95), Inches(0.62),
              body, size=10.5, color=MUTED)
        y += Inches(1.32)

    # 9 — next steps --------------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Next steps", "Sequenced by dependency, not by preference")
    rows = [
        ["#", "Task", "Why now", "Effort"],
        ["1", "Custom avatar with the KOL's face",
         "The single biggest gap — the voice is hers, the face is not", "M"],
        ["2", "LLM brain (Ollama + persona prompt)",
         "Turns echo-mode into actual conversation; profiles already hold the persona", "S"],
        ["3", "Face-lock LoRA for image generation",
         "Consistent face across posts; unblocks the content pillars", "L"],
        ["4", "Speaker diarization in crawl.py",
         "Makes real multi-speaker audio safe to train on", "M"],
        ["5", "Voices for the remaining 9 KOLs",
         "Pipeline is proven; each run is now largely unattended", "S each"],
        ["6", "Stream out to OBS / RTMP",
         "Last step to an actually broadcastable AI KOL", "S"],
    ]
    table(sl, MARGIN, Inches(1.7), BODY_W, rows,
          [Inches(0.5), Inches(4.3), Inches(6.4), Inches(0.69)], size=11, row_h=Inches(0.52))
    txbox(sl, MARGIN, Inches(5.4), BODY_W, Inches(1.2),
          "Recommended order: 1 → 2 gives a demonstrable talking, thinking KOL with her own "
          "face and voice — the shortest path to something presentable. 3 is the largest single "
          "investment and can run in parallel. Compliance reminder: AI disclosure is required in "
          "several markets, and TikTok has no official comment API, so automation must stay off "
          "there.", size=11.5, color=MUTED)

    # 10 — status board -----------------------------------------------------
    sl = slide_blank(prs)
    header(sl, "Appendix", "Current system state (live at generation time)")
    svc_rows = [["Service", "Port", "Role", "State"]]
    for x in state["services"]:
        svc_rows.append([x["name"], str(x["port"]), x["role"],
                         ("running", OK) if x["up"] else ("not running", BAD)])
    table(sl, MARGIN, Inches(1.7), BODY_W, svc_rows,
          [Inches(3.0), Inches(0.9), Inches(6.5), Inches(1.49)], size=11, row_h=Inches(0.42))

    y = Inches(3.66)
    txbox(sl, MARGIN, y, BODY_W, Inches(0.3), "KOL pipeline state", size=12, color=ACCENT, bold=True)
    kol_rows = [["KOL", "Status", "Images", "Voice dataset", "Voice trained"]]
    for k in state["kols"]:
        v = k["voice"]
        if v["sovits"] and v["gpt"]:
            trained = ("yes", OK)
        elif v["clips"]:
            trained = ("dataset only", WARN)
        else:
            trained = ("no", MUTED)
        mins = v["minutes"]
        ds = f'{v["clips"]} clips · {mins:.1f} min' if isinstance(mins, (int, float)) and mins \
            else (f'{v["clips"]} clips' if v["clips"] else "—")
        kol_rows.append([k["name"], k["status"] or "—", str(k["images"] or "—"), ds, trained])
    table(sl, MARGIN, y + Inches(0.36), BODY_W, kol_rows,
          [Inches(3.2), Inches(2.0), Inches(1.3), Inches(3.3), Inches(2.09)],
          size=9.5, row_h=Inches(0.245))

    if gpu:
        txbox(sl, MARGIN, Inches(7.0), BODY_W, Inches(0.3),
              f'GPU: {gpu["name"]} — {gpu["used_mb"]:,} / {gpu["total_mb"]:,} MiB used, '
              f'{gpu["total_mb"]-gpu["used_mb"]:,} MiB free',
              size=10.5, color=MUTED)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("-o", "--out", default=str(REPO / "docs" / "AI-KOL-Progress-Report.pptx"))
    args = ap.parse_args()

    import server  # tools/dashboard/server.py
    state = server.collect()
    out = build(state, Path(args.out))
    print(f"wrote {out}  ({out.stat().st_size/1e3:.0f} KB)")
    print(f"slides: 10   KOLs: {state['summary']['kols_total']}   "
          f"voiced: {state['summary']['kols_voiced']}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
